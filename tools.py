import os
import json
import base64
import requests
import concurrent.futures
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import comtypes.client
from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
from crewai.tools import BaseTool
from pydantic import Field
from typing import ClassVar, Dict, Any
from dotenv import load_dotenv

load_dotenv()

class GoogleSearchTool(BaseTool):
    name: str = "Google Search"
    description: str = "Search the web using Google Search via Serper API for detailed information."

    def _run(self, query: str) -> str:
        url = "https://google.serper.dev/search"
        payload = json.dumps({
            "q": query,
            "num": 5  # Top 5 results
        })
        headers = {
            'X-API-KEY': os.getenv('SERPER_API_KEY'),
            'Content-Type': 'application/json'
        }
        try:
            response = requests.request("POST", url, headers=headers, data=payload)
            response.raise_for_status()
            results = response.json().get('organic', [])
            
            # Format the output for the agent nicely
            formatted_results = f"Search Results for '{query}':\n\n"
            for result in results:
                formatted_results += f"Title: {result.get('title')}\n"
                formatted_results += f"Snippet: {result.get('snippet')}\n"
                formatted_results += f"Link: {result.get('link')}\n\n"
                
            return formatted_results if results else "No significant results found."
        except Exception as e:
            return f"Error searching Google via Serper: {e}"

class PPTXGeneratorTool(BaseTool):
    name: str = "PPTX Generator"
    description: str = "Generates a PowerPoint presentation from a JSON string describing slides."
    
    # Professional color themes
    THEMES: ClassVar[Dict[str, Dict[str, Any]]] = {
        "modern": {
            "primary": RGBColor(41, 128, 185),      # Professional Blue
            "secondary": RGBColor(52, 73, 94),      # Dark Blue-Gray
            "accent": RGBColor(26, 188, 156),       # Teal
            "background": RGBColor(236, 240, 241),  # Light Gray
            "text": RGBColor(44, 62, 80)            # Dark Gray
        },
        "creative": {
            "primary": RGBColor(142, 68, 173),      # Purple
            "secondary": RGBColor(155, 89, 182),    # Light Purple
            "accent": RGBColor(241, 196, 15),       # Golden Yellow
            "background": RGBColor(250, 250, 250),  # Off-White
            "text": RGBColor(52, 73, 94)            # Dark Blue-Gray
        },
        "corporate": {
            "primary": RGBColor(31, 58, 147),       # Navy Blue
            "secondary": RGBColor(66, 103, 178),    # Blue
            "accent": RGBColor(220, 53, 69),        # Professional Red
            "background": RGBColor(248, 249, 250),  # Light Background
            "text": RGBColor(33, 37, 41)            # Near Black
        }
    }

    def _get_theme_for_topic(self, title: str) -> str:
        """Auto-select theme based on topic keywords"""
        title_lower = title.lower()
        
        if any(word in title_lower for word in ['business', 'corporate', 'finance', 'company']):
            return "corporate"
        elif any(word in title_lower for word in ['art', 'design', 'creative', 'innovation']):
            return "creative"
        else:
            return "modern"
    
    def _apply_gradient_background(self, slide, theme_colors):
        """Apply a subtle gradient background to slide"""
        try:
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_angle = 90.0
            
            # Two-color gradient from background to slightly darker
            fill.gradient_stops[0].color.rgb = theme_colors["background"]
            fill.gradient_stops[1].color.rgb = RGBColor(
                max(0, theme_colors["background"].rgb[0] - 15),
                max(0, theme_colors["background"].rgb[1] - 15),
                max(0, theme_colors["background"].rgb[2] - 15)
            )
        except:
            # If gradient fails, use solid background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = theme_colors["background"]
    
    def _style_title(self, title_shape, theme_colors, is_main_title=False):
        """Apply professional styling to title"""
        if not title_shape or not title_shape.has_text_frame:
            return
            
        text_frame = title_shape.text_frame
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER if is_main_title else PP_ALIGN.LEFT
            paragraph.font.size = Pt(54) if is_main_title else Pt(40)
            paragraph.font.bold = True
            paragraph.font.color.rgb = theme_colors["primary"]
            paragraph.font.name = "Calibri"
    
    def _style_subtitle(self, subtitle_shape, theme_colors):
        """Apply styling to subtitle"""
        if not subtitle_shape or not subtitle_shape.has_text_frame:
            return
            
        text_frame = subtitle_shape.text_frame
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = theme_colors["secondary"]
            paragraph.font.name = "Calibri"
    
    def _style_body_text(self, body_shape, theme_colors):
        """Apply styling to body text"""
        if not body_shape or not body_shape.has_text_frame:
            return
            
        text_frame = body_shape.text_frame
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = theme_colors["text"]
            paragraph.font.name = "Calibri"
            paragraph.space_before = Pt(12)
            paragraph.space_after = Pt(12)
            
            # Style bullet points
            if paragraph.level == 0:
                paragraph.font.size = Pt(20)
            else:
                paragraph.font.size = Pt(18)
    
    def _add_accent_shape(self, slide, theme_colors):
        """Add a decorative accent shape to the slide"""
        try:
            # Add a thin accent bar at the top
            left = Inches(0)
            top = Inches(0)
            width = Inches(10)
            height = Inches(0.1)
            
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, width, height
            )
            
            # Style the shape
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme_colors["accent"]
            shape.line.fill.background()
        except:
            pass  # If shape fails, continue without it
    
    def _create_table(self, slide, table_data, theme_colors, position="center"):
        """Create a professionally styled table on the slide"""
        try:
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            
            if not headers or not rows:
                return None
            
            num_cols = len(headers)
            num_rows = len(rows) + 1  # +1 for header row
            
            # Calculate table dimensions and position
            if position == "center":
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(8)
            elif position == "right":
                left = Inches(5.5)
                top = Inches(2)
                width = Inches(4)
            else:  # left or default
                left = Inches(0.5)
                top = Inches(2)
                width = Inches(4.5)
            
            height = Inches(0.5 * num_rows)
            
            # Create table
            table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
            
            # Style header row
            for col_idx, header in enumerate(headers):
                cell = table.rows[0].cells[col_idx]
                cell.text = str(header)
                
                # Header styling
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme_colors["primary"]
                
                # Text styling
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
            
            # Style data rows
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_value in enumerate(row_data):
                    cell = table.rows[row_idx + 1].cells[col_idx]
                    cell.text = str(cell_value)
                    
                    # Alternating row colors for better readability
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(248, 249, 250)  # Light gray
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
                    
                    # Text styling
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(16)
                    paragraph.font.color.rgb = theme_colors["text"]
                    paragraph.font.name = "Calibri"
                    paragraph.alignment = PP_ALIGN.LEFT
            
            return table
        except Exception as e:
            print(f"Error creating table: {e}")
            return None

    def _run(self, json_content: str) -> str:
        try:
            # Clean up JSON string if it contains markdown code blocks
            if "```json" in json_content:
                json_content = json_content.split("```json")[1].split("```")[0].strip()
            elif "```" in json_content:
                json_content = json_content.split("```")[1].split("```")[0].strip()
                
            data = json.loads(json_content)
            prs = Presentation()
            
            # Set slide size to widescreen (16:9)
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Auto-select theme based on topic
            title = data.get("title", "Presentation")
            theme_name = self._get_theme_for_topic(title)
            theme_colors = self.THEMES[theme_name]
            
            # Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Apply gradient background
            self._apply_gradient_background(slide, theme_colors)
            
            # Style title slide
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            subtitle_shape.text = "Generated by AI"
            
            self._style_title(title_shape, theme_colors, is_main_title=True)
            self._style_subtitle(subtitle_shape, theme_colors)
            
            # Add accent shape
            self._add_accent_shape(slide, theme_colors)

            # Content Slides
            for slide_data in data.get("slides", []):
                slide_type = slide_data.get("type", "bullets")
                
                # Choose layout based on slide type
                if slide_type == "table":
                    # Use blank layout for table-only slides
                    blank_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_layout)
                else:
                    # Use bullet layout for bullets or mixed
                    bullet_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(bullet_slide_layout)
                
                # Apply gradient background
                self._apply_gradient_background(slide, theme_colors)
                
                shapes = slide.shapes
                
                # Add title
                if slide_type == "table":
                    # Manually add title for blank layout
                    title_shape = shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
                    )
                    title_frame = title_shape.text_frame
                    title_para = title_frame.paragraphs[0]
                    title_para.text = slide_data.get("title", "Slide")
                    title_para.font.size = Pt(40)
                    title_para.font.bold = True
                    title_para.font.color.rgb = theme_colors["primary"]
                    title_para.font.name = "Calibri"
                else:
                    title_shape = shapes.title
                    title_shape.text = slide_data.get("title", "Slide")
                    self._style_title(title_shape, theme_colors, is_main_title=False)
                
                # Handle different content types
                if slide_type == "bullets":
                    # Standard bullet slide
                    body_shape = shapes.placeholders[1]
                    tf = body_shape.text_frame
                    
                    for point in slide_data.get("bullets", []):
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                    
                    self._style_body_text(body_shape, theme_colors)
                
                elif slide_type == "table":
                    # Table-only slide
                    table_data = slide_data.get("table", {})
                    if table_data:
                        self._create_table(slide, table_data, theme_colors, position="center")
                
                elif slide_type == "mixed":
                    # Mixed: bullets on left, table on right
                    body_shape = shapes.placeholders[1]
                    
                    # Resize body shape to left half
                    body_shape.left = Inches(0.5)
                    body_shape.width = Inches(4.5)
                    
                    # Add bullets
                    tf = body_shape.text_frame
                    for point in slide_data.get("bullets", []):
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                    
                    self._style_body_text(body_shape, theme_colors)
                    
                    # Add table on right
                    table_data = slide_data.get("table", {})
                    if table_data:
                        self._create_table(slide, table_data, theme_colors, position="right")
                
                # Add accent shape
                self._add_accent_shape(slide, theme_colors)

            output_path = os.path.abspath("output/presentation.pptx")
            os.makedirs("output", exist_ok=True)
            prs.save(output_path)
            return f"Presentation saved to {output_path} (Theme: {theme_name})"
        except Exception as e:
            return f"Error generating PPTX: {e}"

class SlideToImageTool(BaseTool):
    name: str = "Slide to Image"
    description: str = "Converts a PPTX file to a folder of images using PowerPoint COM."

    def _run(self, pptx_path: str) -> str:
        try:
            pptx_path = os.path.abspath(pptx_path)
            output_folder = os.path.abspath("output/slides")
            os.makedirs(output_folder, exist_ok=True)

            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            # powerpoint.Visible = 1 # Keep it invisible if possible, or visible if needed
            
            presentation = powerpoint.Presentations.Open(pptx_path)
            
            # Export each slide
            for i, slide in enumerate(presentation.Slides):
                image_path = os.path.join(output_folder, f"slide_{i+1}.jpg")
                slide.Export(image_path, "JPG")
            
            presentation.Close()
            powerpoint.Quit()
            
            return f"Slides exported to {output_folder}"
        except Exception as e:
            return f"Error converting slides to images: {e}"

class TextToSpeechTool(BaseTool):
    name: str = "Text to Speech"
    description: str = "Converts text to speech audio files using Sarvam AI bulbul:v2."

    def _synthesise_one(self, text: str, out_path: str) -> str:
        """Call Sarvam AI TTS for a single chunk and save as .wav"""
        api_key = os.getenv("SARVAM_API")
        url = "https://api.sarvam.ai/text-to-speech"
        headers = {
            "api-subscription-key": api_key,
            "Content-Type": "application/json"
        }
        payload = {
            "inputs": [text[:1500]],          # Sarvam limit: 1500 chars
            "target_language_code": "en-IN",
            "speaker": "anushka",
            "model": "bulbul:v2",
            "enable_preprocessing": True
        }
        try:
            response = requests.post(url, headers=headers, json=payload, timeout=30)
            response.raise_for_status()
            data = response.json()
            audio_b64 = data["audios"][0]
            audio_bytes = base64.b64decode(audio_b64)
            with open(out_path, "wb") as f:
                f.write(audio_bytes)
            return out_path
        except Exception as e:
            # Fallback: write silence as empty wav so pipeline doesn't break
            print(f"[TTS Warning] Sarvam API error for '{out_path}': {e}")
            self._write_silent_wav(out_path, duration_sec=3)
            return out_path

    def _write_silent_wav(self, path: str, duration_sec: int = 3):
        """Write a minimal silent WAV file as a fallback."""
        import struct, wave
        sample_rate = 22050
        n_samples = sample_rate * duration_sec
        with wave.open(path, 'w') as wf:
            wf.setnchannels(1)
            wf.setsampwidth(2)
            wf.setframerate(sample_rate)
            wf.writeframes(struct.pack('<' + 'h' * n_samples, *([0] * n_samples)))

    def _run(self, json_script: str) -> str:
        try:
            # Clean JSON
            if "```json" in json_script:
                json_script = json_script.split("```json")[1].split("```")[0].strip()
            elif "```" in json_script:
                json_script = json_script.split("```")[1].split("```")[0].strip()

            data = json.loads(json_script)
            output_folder = os.path.abspath("output/audio")
            os.makedirs(output_folder, exist_ok=True)

            # Build list of (text, filepath) pairs
            tasks = []
            title_text = f"Welcome to this presentation on {data.get('title', 'the topic')}."
            tasks.append((title_text, os.path.join(output_folder, "audio_0.wav")))

            for i, slide in enumerate(data.get("slides", [])):
                text = slide.get("voiceover", "")
                if not text:
                    text = f"Slide {i+1}: {slide.get('title', 'Next point')}."
                tasks.append((text, os.path.join(output_folder, f"audio_{i+1}.wav")))

            # Generate all audio files IN PARALLEL
            print(f"  [TTS] Generating {len(tasks)} audio clips in parallel via Sarvam AI...")
            with concurrent.futures.ThreadPoolExecutor(max_workers=min(len(tasks), 6)) as executor:
                futures = {executor.submit(self._synthesise_one, text, path): path
                           for text, path in tasks}
                for future in concurrent.futures.as_completed(futures):
                    result = future.result()
                    print(f"  [TTS] ✓ {os.path.basename(result)}")

            return f"Audio files generated in {output_folder} (Sarvam AI bulbul:v2)"
        except Exception as e:
            return f"Error generating audio: {e}"


class VideoCompilerTool(BaseTool):
    name: str = "Video Compiler"
    description: str = "Compiles images and audio into a final video."

    def _run(self, args: str) -> str:
        # args is ignored, we look at standard paths
        try:
            slides_dir = os.path.abspath("output/slides")
            audio_dir = os.path.abspath("output/audio")
            output_video = os.path.abspath("output/final_video.mp4")

            slides = sorted([os.path.join(slides_dir, f) for f in os.listdir(slides_dir) if f.endswith(".jpg")])
            audios = sorted([os.path.join(audio_dir, f) for f in os.listdir(audio_dir) if f.endswith(".wav")])

            if len(slides) != len(audios):
                return f"Error: Mismatch in slides ({len(slides)}) and audio files ({len(audios)})."

            clips = []
            for slide_path, audio_path in zip(slides, audios):
                audio_clip = AudioFileClip(audio_path)
                slide_clip = ImageClip(slide_path).with_duration(audio_clip.duration)
                slide_clip = slide_clip.with_audio(audio_clip)
                clips.append(slide_clip)

            final_clip = concatenate_videoclips(clips)
            final_clip.write_videofile(output_video, fps=24)

            return f"Video created at {output_video}"
        except Exception as e:
            return f"Error compiling video: {e}"

